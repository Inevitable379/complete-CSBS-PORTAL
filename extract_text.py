import os
import json
import logging

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

logging.basicConfig(level=logging.INFO)

base_dir = r'C:\Users\Avi\OneDrive\Desktop\jain\SEM-2'
output_file = 'file_contents_sem2.json'

allowed_exts = {'.pdf', '.docx', '.pptx'}
results = {}

def extract_pdf(file_path):
    if not PdfReader: return "PyPDF2 not installed"
    try:
        reader = PdfReader(file_path)
        if len(reader.pages) > 0:
            text = reader.pages[0].extract_text()
            return text[:500] if text else "No text found on page 1"
    except Exception as e:
        return f"Error: {e}"
    return "Empty PDF"

def extract_docx(file_path):
    if not docx: return "python-docx not installed"
    try:
        doc = docx.Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs[:10]])
        return text[:500]
    except Exception as e:
        return f"Error: {e}"

def extract_pptx(file_path):
    if not pptx: return "python-pptx not installed"
    try:
        prs = pptx.Presentation(file_path)
        text = []
        if len(prs.slides) > 0:
            for shape in prs.slides[0].shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)[:500]
    except Exception as e:
        return f"Error: {e}"

for folder in os.listdir(base_dir):
    folder_path = os.path.join(base_dir, folder)
    if not os.path.isdir(folder_path):
        continue
    
    results[folder] = {}
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if not os.path.isfile(file_path):
            continue
            
        _, ext = os.path.splitext(filename)
        ext = ext.lower()
        
        if ext not in allowed_exts:
            results[folder][filename] = "Skipped (unsupported extension)"
            continue
            
        logging.info(f"Extracting {filename}...")
        if ext == '.pdf':
            text = extract_pdf(file_path)
        elif ext == '.docx':
            text = extract_docx(file_path)
        elif ext == '.pptx':
            text = extract_pptx(file_path)
            
        results[folder][filename] = text

with open(output_file, 'w', encoding='utf-8') as f:
    json.dump(results, f, indent=2)

print(f"Extraction complete! Saved to {output_file}")
